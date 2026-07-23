"""
Build the RepurposeIQ "Three Use Cases" slide deck (curated narrative).

Generates an editable PowerPoint (.pptx) in the project root:
    RepurposeIQ_Use_Cases.pptx

Layout: cover + 3 case slides (one per use case) + closing summary, styled in the
platform's warm "Clinical Cream & Botanical" theme with the Solix logo.

    python build_case_deck.py
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

ROOT = Path(__file__).parent
LOGO = ROOT / "assets" / "solix_logo_transparent.png"
OUT  = ROOT / "RepurposeIQ_Use_Cases.pptx"
DEMO_DATE = "2 July 2026"

# ── Brand palette (from templates/base.html) ────────────────────────────────
CREAM      = RGBColor(0xF6, 0xF5, 0xF0)
CARD       = RGBColor(0xFF, 0xFD, 0xF9)
SUNKEN     = RGBColor(0xEC, 0xEB, 0xE3)
GREEN      = RGBColor(0x2F, 0x6F, 0x5B)
GREEN_DEEP = RGBColor(0x24, 0x55, 0x46)
GOLD       = RGBColor(0xB8, 0x86, 0x0B)
GOLD_DEEP  = RGBColor(0x8A, 0x63, 0x09)
TERRA      = RGBColor(0xB5, 0x55, 0x2F)
INK        = RGBColor(0x1F, 0x2A, 0x24)
INK_2      = RGBColor(0x45, 0x52, 0x4A)
MUTED      = RGBColor(0x6B, 0x77, 0x6E)
DIM        = RGBColor(0x9A, 0xA3, 0x9C)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LINE       = RGBColor(0xDD, 0xDB, 0xD0)

FONT = "Segoe UI"

SW, SH = Inches(13.333), Inches(7.5)


# ── Low-level helpers ───────────────────────────────────────────────────────

def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _no_line(shape):
    shape.line.fill.background()


def _rect(slide, x, y, w, h, fill, rounded=False, radius=0.055, line=None, shadow=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        x, y, w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        _no_line(shp)
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    if rounded:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    shp.shadow.inherit = False
    return shp


def _text(slide, x, y, w, h, runs, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
          space_after=4, line_spacing=1.05):
    """runs: list of paragraphs; each paragraph is a list of (text, size, color, bold) tuples."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (txt, size, color, bold) in para:
            r = p.add_run()
            r.text = txt
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.name = FONT
    return tb


def _bullets(slide, x, y, w, h, items, *, size=11.5, color=INK_2, gap=5,
            bullet_color=GREEN, line_spacing=1.06):
    """items: list of strings, OR (lead, rest) tuples where lead is bolded ink."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        b = p.add_run(); b.text = "•  "
        b.font.size = Pt(size); b.font.color.rgb = bullet_color; b.font.bold = True
        b.font.name = FONT
        if isinstance(item, tuple):
            lead, rest = item
            r1 = p.add_run(); r1.text = lead
            r1.font.size = Pt(size); r1.font.color.rgb = INK; r1.font.bold = True
            r1.font.name = FONT
            if rest:
                r2 = p.add_run(); r2.text = rest
                r2.font.size = Pt(size); r2.font.color.rgb = color; r2.font.bold = False
                r2.font.name = FONT
        else:
            r = p.add_run(); r.text = item
            r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = False
            r.font.name = FONT
    return tb


def _logo(slide, x, y, width=Inches(1.35)):
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), x, y, width=width)


def _section_label(slide, x, y, w, text, color=GREEN):
    _text(slide, x, y, w, Inches(0.3),
          [[(text.upper(), 10.5, color, True)]], space_after=0)


# ── Slide builders ──────────────────────────────────────────────────────────

def cover(prs):
    s = _blank(prs)
    _bg(s, GREEN_DEEP)
    # botanical accent band
    _rect(s, 0, Inches(5.55), SW, Inches(0.06), GOLD)
    _rect(s, Inches(0.9), Inches(1.0), Inches(0.55), Inches(2.9), GOLD)
    if LOGO.exists():
        s.shapes.add_picture(str(LOGO), Inches(0.9), Inches(0.62), width=Inches(2.2))
    _text(s, Inches(0.85), Inches(2.15), Inches(11.5), Inches(2.4), [
        [("RepurposeIQ", 46, WHITE, True)],
        [("Three Repurposing Use Cases", 24, RGBColor(0xE7, 0xEF, 0xEA), False)],
    ], space_after=8)
    _text(s, Inches(0.9), Inches(4.15), Inches(11.5), Inches(1.2), [
        [("01  ", 13, GOLD, True), ("Palbociclib — non-oncology repurposing", 13, RGBColor(0xD9,0xE4,0xDE), False)],
        [("02  ", 13, GOLD, True), ("Erlotinib — dermatology / ophthalmology indications", 13, RGBColor(0xD9,0xE4,0xDE), False)],
        [("03  ", 13, GOLD, True), ("Vitiligo — candidate molecules for repurposing", 13, RGBColor(0xD9,0xE4,0xDE), False)],
    ], space_after=6, line_spacing=1.2)
    _text(s, Inches(0.9), Inches(6.55), Inches(11.5), Inches(0.5),
          [[("Solix  ·  Drug Repurposing Platform  ·  " + DEMO_DATE, 11, RGBColor(0xAF,0xC3,0xBA), False)]])


def case_slide(prs, *, num, title, subtitle, left_head, left_items,
               understood, solution_items, footer, disease=False):
    s = _blank(prs)
    _bg(s, CREAM)

    # header
    _rect(s, 0, 0, Inches(0.16), SH, GREEN)          # left spine
    _section_label(s, Inches(0.6), Inches(0.42), Inches(6), f"CASE {num:02d}")
    _text(s, Inches(0.6), Inches(0.66), Inches(9.6), Inches(1.0), [
        [(title, 26, INK, True)],
        [(subtitle, 13.5, MUTED, False)],
    ], space_after=3)
    _logo(s, Inches(11.55), Inches(0.5))
    _rect(s, Inches(0.6), Inches(1.72), Inches(12.15), Pt(1.4), LINE)

    top = Inches(2.02)
    colh = Inches(4.55)

    # left card — the drug / disease
    lx, lw = Inches(0.6), Inches(4.85)
    _rect(s, lx, top, lw, colh, CARD, rounded=True, line=LINE)
    _rect(s, lx, top, lw, Inches(0.5), SUNKEN, rounded=True)
    _text(s, lx + Inches(0.28), top + Inches(0.09), lw - Inches(0.5), Inches(0.36),
          [[(left_head, 13, GREEN_DEEP, True)]])
    _bullets(s, lx + Inches(0.28), top + Inches(0.66), lw - Inches(0.5), colh - Inches(0.8),
             left_items, size=11, gap=7)

    # right column — two stacked cards
    rx, rw = Inches(5.68), Inches(7.07)
    # understood
    uh = Inches(1.42)
    _rect(s, rx, top, rw, uh, RGBColor(0xF3, 0xEE, 0xE0), rounded=True,
          line=RGBColor(0xE3, 0xD6, 0xB6))
    _text(s, rx + Inches(0.3), top + Inches(0.14), rw - Inches(0.55), Inches(0.3),
          [[("WHAT WE UNDERSTOOD", 10.5, GOLD_DEEP, True)]])
    _text(s, rx + Inches(0.3), top + Inches(0.5), rw - Inches(0.55), uh - Inches(0.55),
          [[(understood, 12.5, INK, False)]], line_spacing=1.12)

    # solution
    sy = top + uh + Inches(0.18)
    sh = colh - uh - Inches(0.18)
    _rect(s, rx, sy, rw, sh, CARD, rounded=True, line=LINE)
    _rect(s, rx, sy, rw, Inches(0.5), GREEN, rounded=True)
    _text(s, rx + Inches(0.3), sy + Inches(0.09), rw - Inches(0.55), Inches(0.36),
          [[("HOW REPURPOSEIQ SOLVES IT", 12, WHITE, True)]])
    _bullets(s, rx + Inches(0.3), sy + Inches(0.68), rw - Inches(0.6), sh - Inches(0.8),
             solution_items, size=11, gap=6, bullet_color=GREEN)

    # footer
    _text(s, Inches(0.6), Inches(6.95), Inches(12.15), Inches(0.4),
          [[(footer, 9, DIM, False)]])


def closing(prs):
    s = _blank(prs)
    _bg(s, CREAM)
    _rect(s, 0, 0, SW, Inches(1.55), GREEN_DEEP)
    if LOGO.exists():
        s.shapes.add_picture(str(LOGO), Inches(11.35), Inches(0.5), width=Inches(1.5))
    _text(s, Inches(0.6), Inches(0.42), Inches(10), Inches(1.0), [
        [("How the platform works", 26, WHITE, True)],
        [("One engine set behind all three use cases", 13, RGBColor(0xD9,0xE4,0xDE), False)],
    ], space_after=3)

    cards = [
        ("Two directions",
         ["Reverse: drug → new indications (Cases 1 & 2)",
          "Forward / pathway-first: disease → candidate molecules (Case 3)"]),
        ("Every claim is a scored pair",
         ["Drug ↔ indication scored by the canonical 6-dimension repurposing score",
          "Direction-aware mechanism check (corrector vs mimic)"]),
        ("Decision-grade enrichment",
         ["Probability-of-success + phase-aware regulatory verdict",
          "505(b)(2) framing; optional docking / QM on leads"]),
        ("Grounded in real data",
         ["ChEMBL mechanisms, Open Targets associations, curated pathways, MeSH",
          "Understands partial fits: primary disorder ↔ systemic manifestations"]),
    ]
    x0, y0 = Inches(0.6), Inches(1.95)
    cw, ch = Inches(5.95), Inches(2.15)
    gap = Inches(0.25)
    for i, (head, items) in enumerate(cards):
        cx = x0 + (cw + gap) * (i % 2)
        cy = y0 + (ch + gap) * (i // 2)
        _rect(s, cx, cy, cw, ch, CARD, rounded=True, line=LINE)
        _rect(s, cx, cy + Inches(0.2), Inches(0.09), ch - Inches(0.4), GREEN)
        _text(s, cx + Inches(0.32), cy + Inches(0.2), cw - Inches(0.5), Inches(0.4),
              [[(head, 14, GREEN_DEEP, True)]])
        _bullets(s, cx + Inches(0.32), cy + Inches(0.72), cw - Inches(0.55), ch - Inches(0.85),
                 items, size=11, gap=5)

    _text(s, Inches(0.6), Inches(6.98), Inches(12), Inches(0.4),
          [[("RepurposeIQ  ·  Solix  ·  " + DEMO_DATE, 9.5, DIM, False)]])


# ── Content ─────────────────────────────────────────────────────────────────

def build():
    prs = Presentation()
    prs.slide_width, prs.slide_height = SW, SH

    cover(prs)

    # CASE 1 — Palbociclib
    case_slide(
        prs, num=1,
        title="Palbociclib — Non-Oncology Repurposing",
        subtitle="Finding new indications outside cancer for a CDK4/6 inhibitor",
        left_head="THE DRUG",
        left_items=[
            ("Class:  ", "selective, oral CDK4/6 inhibitor (Ibrance)"),
            ("Target / MoA:  ", "inhibits CDK4 & CDK6 → blocks Rb phosphorylation → halts E2F-driven G1→S cell-cycle progression"),
            ("Approved use:  ", "HR+ / HER2− advanced or metastatic breast cancer, with endocrine therapy"),
            ("Key pathways:  ", "CDK4/6–Cyclin D–Rb–E2F; cellular senescence (SASP); NF-κB inflammatory signaling"),
        ],
        understood="“Non-oncology repurposing” = surface NEW disease areas outside cancer that exploit "
                   "CDK4/6 and senescence / inflammation biology — not another tumour type.",
        solution_items=[
            ("Engine:  ", "reverse-repurposing (drug → indications) with the therapeutic-area filter set to EXCLUDE oncology"),
            ("Method:  ", "drug targets (CDK4/6) → Open Targets target→disease → subtract known uses → canonical 6-dimension score"),
            ("Mechanistic angle:  ", "senescence & NF-κB rationale points at fibrotic and immune-inflammatory directions"),
            ("Each candidate carries:  ", "canonical repurposing score, probability-of-success, phase-aware regulatory verdict, 505(b)(2) angle"),
        ],
        footer="Reverse-repurposing engine · ChEMBL mechanisms + Open Targets associations · scored by the platform's canonical drug↔indication pair score.",
    )

    # CASE 2 — Erlotinib
    case_slide(
        prs, num=2,
        title="Erlotinib — Dermatology / Ophthalmology",
        subtitle="Repurposing an EGFR inhibitor into skin & eye indications",
        left_head="THE DRUG",
        left_items=[
            ("Class:  ", "reversible, oral EGFR (HER1) tyrosine-kinase inhibitor (Tarceva)"),
            ("Target / MoA:  ", "competes at the EGFR ATP site → blocks receptor autophosphorylation → shuts down RAS–MAPK and PI3K–AKT signaling"),
            ("Approved use:  ", "EGFR-mutation NSCLC; locally advanced / metastatic pancreatic cancer (with gemcitabine)"),
            ("Key pathways:  ", "EGFR–RAS–MAPK; EGFR–PI3K–AKT; keratinocyte proliferation & ocular-surface homeostasis"),
        ],
        understood="Repurpose Erlotinib into SKIN and EYE indications where EGFR-driven hyperproliferation "
                   "or inflammation is the mechanism — a targeted-area repurposing, not oncology.",
        solution_items=[
            ("Engine:  ", "reverse-repurposing with therapeutic-area filter = Dermatology + Ophthalmology"),
            ("Area logic:  ", "uses Open Targets' own therapeuticAreas annotation, not a hand-kept keyword list"),
            ("Method:  ", "EGFR target → area-filtered disease associations → canonical score + direction-aware mechanism check"),
            ("Built-in context:  ", "EGFR inhibition's known cutaneous pharmacology is surfaced as mechanistic support and a safety consideration"),
        ],
        footer="Reverse-repurposing engine, therapeutic-area scoped · Open Targets therapeuticAreas + ChEMBL · canonical pair score with direction-aware mechanism.",
    )

    # CASE 3 — Vitiligo
    case_slide(
        prs, num=3,
        title="Vitiligo — Candidate Molecules",
        subtitle="Disease → drug: finding molecules to repurpose for vitiligo",
        left_head="THE DISEASE",
        left_items=[
            ("Definition:  ", "chronic autoimmune depigmentation — progressive loss of epidermal melanocytes → white macules"),
            ("Epidemiology:  ", "~0.5–1% global prevalence; high psychosocial burden"),
            ("Core biology:  ", "CD8+ T-cell melanocyte killing; IFN-γ → JAK1/JAK2–STAT1 → CXCL9/10/11 → CXCR3+ T-cell recruitment; melanocyte oxidative stress"),
            ("Treatment gap:  ", "steroids / calcineurin inhibitors, phototherapy; topical ruxolitinib (JAK1/2) recently approved — validates the axis"),
        ],
        understood="Disease → drug direction: find CANDIDATE MOLECULES to repurpose FOR vitiligo, "
                   "mechanistically anchored on its driver pathway.",
        solution_items=[
            ("Engines:  ", "forward discovery (disease → compound) + pathway-first screen anchored on JAK-STAT / IFN-γ signaling"),
            ("Method:  ", "vitiligo → disease genes & pathways → drugs that modulate them in the corrective direction → canonical pair score"),
            ("Illustrative class:  ", "JAK inhibitors (ruxolitinib, tofacitinib, baricitinib) and other IFN-γ / CXCR3-axis modulators"),
            ("Each candidate:  ", "mechanistic evidence chain, canonical score, probability-of-success; optional docking / QM on leads"),
        ],
        footer="Forward discovery + pathway-first engines (JAK-STAT anchor) · disease resolved via MeSH + Open Targets · canonical pair score.",
        disease=True,
    )

    closing(prs)

    prs.save(str(OUT))
    return OUT


if __name__ == "__main__":
    out = build()
    print(f"Saved deck: {out}  ({out.stat().st_size // 1024} KB, {5} slides)")
