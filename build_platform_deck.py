"""
Build the RepurposeIQ platform architecture deck (with per-slide speaker notes).
Run:  .venv312\\Scripts\\python.exe build_platform_deck.py
Output: RepurposeIQ_Platform_Architecture.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── Theme (platform "Clinical Cream & Botanical") ──────────────────────────────
GREEN  = RGBColor(0x2F, 0x6F, 0x5B)
CREAM  = RGBColor(0xFB, 0xFA, 0xF6)
CREAM2 = RGBColor(0xF3, 0xF1, 0xE9)
INK    = RGBColor(0x1F, 0x2A, 0x24)
MUTE   = RGBColor(0x45, 0x52, 0x4A)
GOLD   = RGBColor(0xB8, 0x86, 0x0B)
LINE   = RGBColor(0xD8, 0xD4, 0xC8)
WHITEISH = RGBColor(0xFF, 0xFD, 0xF9)

W, H = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]


def _bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _box(slide, l, t, w, h, fill=None, line=None, line_w=1.0):
    from pptx.enum.shapes import MSO_SHAPE
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sp.adjustments[0] = 0.06
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def _text(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
          space_after=4, line_spacing=1.0):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    first = True
    for item in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align; p.space_after = Pt(space_after); p.line_spacing = line_spacing
        p.space_before = Pt(0)
        segs = item.get("segs") or [(item.get("text", ""), item)]
        for text, st in segs:
            r = p.add_run(); r.text = text
            f = r.font
            f.size = Pt(st.get("size", 16)); f.name = "Calibri"
            f.bold = st.get("bold", False); f.italic = st.get("italic", False)
            f.color.rgb = st.get("color", INK)
        if item.get("bullet"):
            _set_bullet(p, item.get("bullet_color", GREEN))
        p.level = item.get("level", 0)
    return tb


def _set_bullet(p, color):
    # add a simple round bullet char via XML
    from pptx.oxml.ns import qn
    pPr = p._pPr if p._pPr is not None else p.get_or_add_pPr()
    for tag in ("a:buNone", "a:buChar", "a:buAutoNum"):
        for e in pPr.findall(qn(tag)):
            pPr.remove(e)
    buFont = pPr.makeelement(qn("a:buFont"), {"typeface": "Arial"})
    buChar = pPr.makeelement(qn("a:buChar"), {"char": "▪"})
    pPr.append(buFont); pPr.append(buChar)


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def title_slide(kicker, title, subtitle, note):
    s = prs.slides.add_slide(BLANK); _bg(s, GREEN)
    _box(s, Inches(0), Inches(0), Inches(0.32), H, fill=GOLD)
    _text(s, Inches(0.9), Inches(2.2), Inches(11.5), Inches(0.5),
          [{"text": kicker, "size": 18, "bold": True, "color": RGBColor(0xCF,0xE0,0xD8)}])
    _text(s, Inches(0.9), Inches(2.7), Inches(11.5), Inches(1.8),
          [{"text": title, "size": 46, "bold": True, "color": WHITEISH}])
    _text(s, Inches(0.9), Inches(4.6), Inches(11.0), Inches(1.2),
          [{"text": subtitle, "size": 20, "color": RGBColor(0xDC,0xE7,0xE1)}], line_spacing=1.15)
    notes(s, note)
    return s


def section_slide(num, title, note):
    s = prs.slides.add_slide(BLANK); _bg(s, INK)
    _text(s, Inches(0.9), Inches(2.9), Inches(2), Inches(1),
          [{"text": num, "size": 60, "bold": True, "color": GOLD}])
    _text(s, Inches(0.9), Inches(4.0), Inches(11.5), Inches(1.4),
          [{"text": title, "size": 34, "bold": True, "color": WHITEISH}])
    notes(s, note)
    return s


def content_slide(title, subtitle, bullets, note, two_col=None, footer=None):
    """bullets: list of (text, level, bold?) tuples. two_col: (left_bullets, right_bullets)."""
    s = prs.slides.add_slide(BLANK); _bg(s, CREAM)
    _box(s, Inches(0.0), Inches(0.0), W, Inches(1.15), fill=GREEN)
    _text(s, Inches(0.6), Inches(0.18), Inches(12.2), Inches(0.6),
          [{"text": title, "size": 26, "bold": True, "color": WHITEISH}])
    if subtitle:
        _text(s, Inches(0.62), Inches(0.75), Inches(12.2), Inches(0.35),
              [{"text": subtitle, "size": 13, "color": RGBColor(0xCF,0xE0,0xD8)}])

    def _render(bl, l, t, w, h):
        runs = []
        for b in bl:
            txt, lvl, bold = (b + (False,))[:3] if isinstance(b, tuple) else (b, 0, False)
            runs.append({"text": txt, "size": 15 if lvl == 0 else 13.5,
                         "bold": bold, "color": INK if lvl == 0 else MUTE,
                         "bullet": True, "bullet_color": GREEN if lvl == 0 else GOLD,
                         "level": lvl, "space_after": 7 if lvl == 0 else 4})
        _text(s, l, t, w, h, runs, line_spacing=1.05)

    if two_col:
        _render(two_col[0], Inches(0.6), Inches(1.45), Inches(6.1), Inches(5.6))
        _box(s, Inches(6.75), Inches(1.5), Pt(1), Inches(5.3), fill=LINE)
        _render(two_col[1], Inches(6.95), Inches(1.45), Inches(6.0), Inches(5.6))
    else:
        _render(bullets, Inches(0.6), Inches(1.45), Inches(12.1), Inches(5.6))

    if footer:
        _text(s, Inches(0.6), Inches(7.02), Inches(12.1), Inches(0.35),
              [{"text": footer, "size": 11, "italic": True, "color": GOLD}])
    notes(s, note)
    return s


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDES
# ═══════════════════════════════════════════════════════════════════════════════

title_slide(
    "CLINICAL REPURPOSING OS",
    "RepurposeIQ",
    "Platform, data pipeline and system architecture — an evidence-first, honestly-validated "
    "drug-repurposing engine for 505(b)(2) decision-making.",
    "Opening slide. RepurposeIQ is an internal decision-support platform for drug repurposing, "
    "built for Solix / the Alembic 505(b)(2) engagement. Frame the talk: we will walk the whole "
    "system top to bottom — the problem, the architecture, the data sources, every functional "
    "module, the scoring science, the docking stack, and — most importantly — what we have "
    "rigorously validated versus what remains heuristic. The through-line of this deck is HONESTY: "
    "the platform is designed so a pharma-grade, skeptical audience can trust the numbers, because "
    "every claim is either validated against a gold standard or explicitly labelled as an estimate.")

content_slide(
    "What RepurposeIQ Is", "One platform, three questions",
    [
        ("Forward: given a DISEASE, which existing drugs could treat it? (Compound Discovery)", 0, True),
        ("Reverse: given a DRUG you already hold approval for, what NEW indications exist? — the 505(b)(2) search (Repurpose Molecule)", 0, True),
        ("Pathway / target-first: enter through a biological PATHWAY or a novel TARGET and reach drug↔indication pairs", 0, True),
        ("Every hypothesis is scored on real mechanistic + clinical evidence, wrapped in a 505(b)(2) feasibility and regulatory view", 0, False),
        ("Backed by real structure-based docking, quantum chemistry, PBPK pharmacokinetics and a validated knowledge-graph predictor", 0, False),
        ("Design stance: disease-AGNOSTIC (no hardcoded disease maps), evidence-first, fail-soft, and validated at the data layer", 0, False),
    ],
    "RepurposeIQ answers three linked questions with one shared scoring core. The FORWARD direction "
    "is classic discovery (disease → drugs). The REVERSE direction is the commercial heart of the "
    "Alembic engagement: take a molecule the client already has approval for and find new indications "
    "that can be pursued via the FDA 505(b)(2) pathway (which lets you rely on existing safety data). "
    "The PATHWAY and NOVEL-TARGET entries let a scientist start from biology rather than a disease "
    "name. Emphasise the design stance — no disease is hardcoded, so the same engine works for "
    "neurology, oncology, cardiology, etc. Everything is fail-soft: if a data source is down, the "
    "platform degrades gracefully rather than breaking.")

content_slide(
    "The Problem & the 505(b)(2) Opportunity", "Why this platform exists",
    None,
    "Left: the scientific and commercial problem. Right: the platform's answer. The key commercial "
    "insight is the 505(b)(2) regulatory pathway — a repurposed molecule can lean on existing safety "
    "and manufacturing data, dramatically cutting cost and time versus a new molecular entity. The "
    "hard part is FINDING and RANKING credible drug-disease pairs from a sea of weak associations, "
    "and doing it with evidence a pharma partner will actually trust. The platform exists because "
    "generic association scores are noisy and un-trustworthy; RepurposeIQ layers real mechanistic, "
    "clinical, structural and regulatory evidence and — critically — tells you how much to trust each number.",
    two_col=(
        [("The problem", 0, True),
         ("Repurposing hypotheses are cheap to generate, hard to rank credibly", 1),
         ("Genetic-association scores are noisy; most links are weak/coincidental", 1),
         ("Pharma partners distrust un-validated 'AI scores'", 1),
         ("Success ≠ plausibility — most plausible pairs still fail in the clinic", 1)],
        [("The opportunity", 0, True),
         ("505(b)(2): reuse existing safety/CMC data → faster, cheaper approval", 1),
         ("A known-safe molecule + a new indication = de-risked asset", 1),
         ("Need: rank pairs on REAL evidence + honest confidence", 1),
         ("Need: regulatory + IP + portfolio fit in one view", 1)],
    ))

section_slide("01", "System & Technology Architecture",
    "Section divider. Next we cover the layered architecture, the technology stack, the compute "
    "environments, and the full data-source map. This is the 'how it's built' portion before we "
    "go module by module.")

content_slide(
    "System Architecture — Layered", "Presentation → Services → Data → Compute",
    [
        ("PRESENTATION — Flask app (flask_app.py), Jinja/HTML templates, 3Dmol.js viewer, server-rendered pages + JSON APIs", 0, True),
        ("SERVICES — ~50 Python modules: resolution, scoring engines, docking, QM, PBPK, KG, regulatory, dossier", 0, True),
        ("DATA — PostgreSQL (chembl_33 + neurorepurpose/Hetionet), local KG files (DRKG, RepoDB), disk caches; ~20 live external APIs", 0, True),
        ("COMPUTE — AutoDock Vina (native), GFN2-xTB quantum engine, scikit-learn models, scipy ODE PBPK, RDKit cheminformatics", 0, True),
        ("CROSS-CUTTING — centralized config, shared HTTP client (retry/backoff), disk cache, fail-soft error handling everywhere", 0, False),
    ],
    "Four layers. PRESENTATION is a Flask application serving both server-rendered pages and a JSON "
    "API consumed by the front-end JavaScript; molecular structures render in-browser with 3Dmol.js. "
    "SERVICES is the bulk of the IP — around fifty focused modules. DATA spans a local PostgreSQL "
    "(the validated ChEMBL 33 database plus a neurorepurpose DB holding the Hetionet knowledge "
    "graph), local knowledge-graph files, and roughly twenty live external APIs. COMPUTE is the heavy "
    "science: real docking, real quantum chemistry, fitted ML models, and ODE-based pharmacokinetics. "
    "The whole thing is glued by cross-cutting infrastructure — one config source of truth, one HTTP "
    "client with retry/backoff, a shared disk cache, and fail-soft behaviour so no single dependency "
    "can take the platform down.")

content_slide(
    "Technology Stack", "What it runs on",
    None,
    "Walk the stack. The app is Flask on Python 3.12 in the .venv312 virtualenv. Data lives in "
    "PostgreSQL on port 5433 — two databases: chembl_33 (the validated ChEMBL release, shared with "
    "the POZ project) and neurorepurpose (Hetionet + platform tables). The heavy-science compute uses "
    "separate micromamba environments under .qc/ so their dependencies never collide with the app. "
    "Cheminformatics is RDKit; ML is scikit-learn; the PBPK solver is scipy LSODA. Docking is the "
    "official AutoDock Vina Windows binary driven through a drugdisc-agent env (Meeko, pdbfixer, "
    "OpenBabel). Quantum chemistry is GFN2-xTB via its own env. Everything is Windows-native.",
    two_col=(
        [("Application", 0, True),
         ("Python 3.12 (.venv312), Flask, Jinja2 templates", 1),
         ("3Dmol.js in-browser molecular viewer", 1),
         ("psycopg2 connection pools", 1),
         ("Data / storage", 0, True),
         ("PostgreSQL :5433 — chembl_33 + neurorepurpose", 1),
         ("Local KG files (DRKG, RepoDB, Hetionet edges)", 1),
         ("Disk caches (services/diskcache.py) + JSON caches", 1)],
        [("Scientific compute", 0, True),
         ("RDKit — cheminformatics, descriptors, conformers", 1),
         ("scikit-learn — DWPC predictor, PoS model", 1),
         ("scipy — LSODA ODE solver (PBPK)", 1),
         ("AutoDock Vina 1.2.7 (native) + Meeko/pdbfixer/OpenBabel", 1),
         ("GFN2-xTB — semi-empirical quantum chemistry", 1),
         ("numpy / scipy.sparse — DWPC metapath matrices", 1)],
    ))

content_slide(
    "Compute Environments", "Isolated by design",
    [
        (".venv312 — the Flask app + core services (RDKit, scikit-learn, scipy, psycopg2)", 0, True),
        (".qc/ micromamba root — hosts two isolated envs so heavy deps never clash with the app:", 0, True),
        ("  • qc — GFN2-xTB quantum chemistry engine (services/qc_engine.py)", 1),
        ("  • drugdisc-agent — docking chemistry: RDKit, Meeko, pdbfixer, OpenBabel, posebusters", 1),
        ("vendor/vina_bin/vina.exe — official AutoDock Vina 1.2.7 Windows binary (the engine)", 0, True),
        ("vendor/AtomisticSkills — vendored open drug-discovery skills (MIT lab) used as the reference recipes", 0, False),
        ("App shells out to env pythons via subprocess (worker scripts) → clean dependency isolation", 0, False),
    ],
    "A key engineering decision: the heavy scientific dependencies live in SEPARATE micromamba "
    "environments, not the app venv. AutoDock Vina and fpocket have no Windows conda build, so we use "
    "the official Vina Windows CLI binary and run the prep chemistry (Meeko, pdbfixer, OpenBabel) in a "
    "dedicated drugdisc-agent env. Quantum chemistry runs in its own qc env. The Flask app never "
    "imports these heavy libraries directly — it writes a job file and shells out to a worker script "
    "running under the right env python, then reads back JSON. This keeps the app lightweight and "
    "means a broken scientific dependency can never crash the web tier. AtomisticSkills (from MIT's "
    "Learning Matter lab) is vendored as the reference implementation for the docking recipes.")

content_slide(
    "Data Sources — Curated & Local", "The trusted foundation",
    None,
    "RepurposeIQ separates data into curated/local (trusted, validated once, queried fast) and live "
    "external APIs (fresh, but rate-limited and cached). This slide is the local foundation. ChEMBL 33 "
    "is the backbone — compounds, targets, mechanisms, indications, activities — restored and "
    "integrity-checked (the data-validation workstream). Hetionet is the curated biomedical knowledge "
    "graph (~47k nodes, gold-standard 'treats' edges) that powers the validated predictor. DRKG and "
    "RepoDB add drug-repurposing edges and a labelled outcome benchmark. MeSH and the Orange Book add "
    "disease vocabulary and FDA patent/exclusivity status. The point for a pharma audience: the "
    "foundation is authoritative, validated once at the data layer, not scraped ad hoc.",
    two_col=(
        [("ChEMBL 33 (PostgreSQL)", 0, True),
         ("Compounds, targets, mechanisms, indications, activities (pChEMBL)", 1),
         ("Restore-verified, integrity-checked", 1),
         ("Hetionet v1.0 (in neurorepurpose DB)", 0, True),
         ("~47k nodes, typed edges; 755 curated Compound-treats-Disease", 1),
         ("Powers the DWPC plausibility predictor", 1)],
        [("DRKG — Drug Repurposing KG", 0, True),
         ("Neurological subset; drug↔gene↔disease triples", 1),
         ("RepoDB (data/external/repodb)", 0, True),
         ("13,558 labelled pairs: 8,931 Approved / 4,627 Failed", 1),
         ("The outcome gold standard for validation", 1),
         ("MeSH + FDA Orange Book", 0, True),
         ("Disease vocabulary; patent & exclusivity status", 1)],
    ))

content_slide(
    "Data Sources — Live External APIs", "~20 authoritative endpoints, cached & fail-soft",
    None,
    "The live layer. Each call goes through the shared HTTP client (retry, backoff, timeout) and is "
    "disk-cached, so we are polite to upstreams and resilient. Open Targets supplies genetic "
    "target-disease associations and pathways; STRING supplies the protein-protein interaction "
    "network; Reactome supplies pathway membership. Clinical reality comes from ClinicalTrials.gov "
    "(plus WHO ICTRP and India's CTRI for global trials), PubMed for literature co-mention, and "
    "openFDA for adverse-event (FAERS) and drug-label signals. Structure comes from RCSB PDB, "
    "AlphaFold and UniProt. Regulatory landscape pulls FDA, EMA, MHRA and Indian sources plus Google "
    "Patents. Stress that all of this is cached and fail-soft — the platform never blocks on a slow API.",
    two_col=(
        [("Biology & mechanism", 0, True),
         ("Open Targets — genetic target↔disease associations, pathways", 1),
         ("STRING — protein-protein interaction network", 1),
         ("Reactome — pathway membership", 1),
         ("UniProt — gene→protein identity resolution", 1),
         ("Clinical & safety", 0, True),
         ("ClinicalTrials.gov, WHO ICTRP, CTRI (India) — trials", 1),
         ("PubMed / NCBI e-utils — literature co-mention", 1),
         ("openFDA — FAERS adverse events + drug labels", 1)],
        [("Structure", 0, True),
         ("RCSB PDB — experimental structures", 1),
         ("AlphaFold (EBI) — predicted models (pLDDT-filtered)", 1),
         ("Regulatory & IP", 0, True),
         ("FDA (Orange Book, labels, approvals)", 1),
         ("EMA, MHRA (UK), CTRI/CDSCO (India)", 1),
         ("Google Patents — patent landscape", 1),
         ("All via services/http_client.py + disk cache (fail-soft)", 1)],
    ))

content_slide(
    "End-to-End Data Pipeline", "From a query to a ranked, dossier-backed hypothesis",
    [
        ("1. RESOLVE — disease → ontology (MeSH/EFO/MONDO) + genes/pathways; drug → ChEMBL id + targets + SMILES", 0, True),
        ("2. GENERATE candidates — DB indications, ChEMBL, Open Targets associations, Hetionet novelty, pathway/target entry", 0, True),
        ("3. GATHER evidence — target overlap, pathways, PPI, trials, literature, FAERS safety, directional KG triples", 0, True),
        ("4. SCORE — canonical 6-dimension composite (renormalized) × guardrail multipliers (collapsed)", 0, True),
        ("5. CALIBRATE — enrichment vs random-pair null + validated DWPC plausibility P(treats) where covered", 0, True),
        ("6. ENRICH — PoS, regulatory verdict, developability, docking, PBPK, evidence dossier, portfolio fit", 0, True),
        ("7. PRESENT — ranked cards, confidence tiers, honest labels, 505(b)(2) dossier, knowledge graph", 0, True),
    ],
    "This is the spine of the platform — every page is a view onto this pipeline. Resolution turns a "
    "free-text query into structured identifiers and biology. Candidate generation pulls from multiple "
    "orthogonal sources so we discover rather than merely confirm. Evidence gathering assembles the "
    "raw signals. Scoring combines them in the canonical composite. Calibration makes the number "
    "interpretable — both relative to a random-pair background and, where the pair is in the "
    "knowledge graph, via the validated plausibility probability. Enrichment adds the decision "
    "context a pharma team needs — probability of success, regulatory and IP status, developability, "
    "3D docking, pharmacokinetics. Presentation renders it with honest confidence labels. Every step "
    "is cached and fail-soft.")

content_slide(
    "Platform Surfaces (Pages)", "Nine working views onto the engine",
    None,
    "The user-facing surfaces, all sharing the same scoring core so a number means the same thing "
    "everywhere. Compound Discovery is the forward disease→drug board. Repurpose Molecule is the "
    "reverse drug→indication 505(b)(2) search. Pathway Screen and Novel Targets are the biology-first "
    "entries. Knowledge Graph visualises the evidence chain. The Analysis / Dossier page is the deep "
    "per-candidate view with docking, QM, PBPK and the 505(b)(2) write-up. Database, Validation, "
    "Business Case, Docs and Settings round it out. Nav is consistent; each page is a thin front-end "
    "over the JSON API.",
    two_col=(
        [("Discovery & repurposing", 0, True),
         ("Compound Discovery (/discover) — disease → drugs", 1),
         ("Repurpose Molecule (/repurpose) — drug → indications, 505(b)(2)", 1),
         ("Pathway Screen (/pathways) — pathway → drug↔indication", 1),
         ("Novel Targets (/novel-targets) — inferred targets → drugs", 1),
         ("Analysis / Dossier (/analysis) — per-candidate deep dive", 1)],
        [("Evidence & context", 0, True),
         ("Knowledge Graph (/graph) — evidence-chain visualisation", 1),
         ("Database (/database) — browse the validated data", 1),
         ("Validation (/validation) — data-quality narrative", 1),
         ("Business Case (/business-case) — 505(b)(2) economics", 1),
         ("Docs & Settings (/docs, /settings)", 1)],
    ))

section_slide("02", "Resolution & Discovery Engines",
    "Section divider. Now module-by-module through the engines that turn a query into scored "
    "candidates: disease and drug resolution, the forward 5-screen engine, reverse repurposing, "
    "pathway-first, and novel targets.")

content_slide(
    "Disease Resolution Layer", "Keeping the platform disease-agnostic",
    [
        ("disease_ontology.py — resolves a disease to Open Targets (EFO/MONDO), pulls genes (genetic weights), pathways, PPI", 0, True),
        ("disease_id.py — ID reconciliation across ChEMBL MeSH, OT MONDO/EFO, RepoDB/UMLS; canonical keys", 0, True),
        ("disease_resolver.py — free-text / alias / abbreviation → canonical heading (handles 'PD', 'alzheimers')", 0, True),
        ("therapeutic_context.py — resolves therapeutic area / route so scoring stays disease-agnostic (no hardcoded maps)", 0, True),
        ("disease_config.py — per-disease configuration hooks", 0, False),
        ("Manifestation-aware: 'systemic sclerosis with pulmonary manifestations' → primary disorder + partial-fit flag", 0, False),
    ],
    "Resolution is where disease-agnosticism is enforced. A user types anything — an abbreviation, an "
    "alias, a manifestation-qualified phrase — and this layer maps it to a canonical ontology term, "
    "then fetches the disease's genetic gene set (with association weights), pathways and PPI "
    "neighbourhood from Open Targets and STRING. The ID reconciliation module is important because "
    "the same disease has different identifiers in ChEMBL (MeSH), Open Targets (MONDO/EFO) and RepoDB "
    "(UMLS); we canonicalise so joins work. Nothing about specific diseases is hardcoded — the same "
    "code path serves any therapy area, which is what makes the platform generalise.")

content_slide(
    "Drug / Compound Resolution", "From a name to a molecule",
    [
        ("compound_validator.py — validates a molecule, RDKit sanity, canonical SMILES, identity checks", 0, True),
        ("neuro_db_service.py — the platform database service over chembl_33 (compounds, targets, mechanisms, activities)", 0, True),
        ("Resolution: name / ChEMBL id → compound row → protein targets (genes), max_phase, known indications, SMILES", 0, True),
        ("Pharmacologic-class awareness: distinguishes a drug from a class label (e.g. 'CDK4/6 inhibitor')", 0, False),
        ("Honest failure: unresolved molecule → explicit error, never docks/scores the wrong compound", 0, False),
    ],
    "The drug side of resolution. Given a name or ChEMBL id, we pull the compound's protein targets "
    "(as gene symbols), its global development phase, its known indications and its SMILES from the "
    "validated ChEMBL database. The compound validator does RDKit-level sanity and identity checks so "
    "we never silently score or dock the wrong molecule — if a name can't be resolved, the platform "
    "says so rather than guessing. It also distinguishes a specific drug from a pharmacologic-class "
    "label, which matters when trials list a class as the 'condition'.")

content_slide(
    "Forward Engine — 5-Screen Repurposing", "repurposing_engine.py",
    [
        ("Given a disease, builds a candidate pool: local DB indications + ChEMBL API + Hetionet-novelty compounds", 0, True),
        ("Hetionet novelty: Compound→Gene←Disease shared-gene links surface drugs NOT already indicated — true discovery", 0, True),
        ("Scores every candidate with the canonical score_compound_for_disease() — the shared scoring core", 0, True),
        ("repurposing_scorer.py — a complementary DB-side scorer (indication / target / activity / network signals)", 0, True),
        ("Augments thin results, resolves missing targets in batch, direction-aware pathway scoring", 0, False),
        ("Output: ranked candidates with full score breakdown, novelty flags, and disease context", 0, False),
    ],
    "The forward engine turns a disease into a ranked drug list. Its cleverness is candidate "
    "generation from multiple orthogonal sources — the local database, the ChEMBL API, and crucially "
    "the Hetionet knowledge graph, which surfaces drugs connected to the disease through shared genes "
    "even when they were never indicated for it. That is what makes the engine DISCOVER new "
    "hypotheses rather than just confirm known ones. Every candidate then flows through the single "
    "canonical scoring function shared by all surfaces, so the number is consistent everywhere. A "
    "second, database-centric scorer contributes indication/target/activity/network signals when "
    "useful.")

content_slide(
    "The Canonical Pair Score", "score_compound_for_disease() — six evidence dimensions",
    None,
    "This is THE scoring function — used identically by every surface via canonical_pair_score(), so "
    "a (drug, disease) pair scores the same on the repurpose card, the dossier and the pathway screen. "
    "Six weighted dimensions. TARGET is genetics-weighted overlap of the drug's targets with the "
    "disease's associated genes. PATHWAY and PPI capture network cohesion, degree-damped so hub genes "
    "don't dominate. CLINICAL is real trials + literature for THIS pair (not global phase — that "
    "would be reverse causation). INDICATION and REGULATORY capture existing-use and development "
    "status. On top sit bounded bonuses for knowledge-graph network evidence and directional "
    "literature triples. The weights shown are the base; the next slide explains the 2026 "
    "renormalization that fixed systematic harshness.",
    two_col=(
        [("Weighted dimensions (base)", 0, True),
         ("target 0.25 — genetics-weighted target overlap", 1),
         ("pathway 0.20 — Reactome pathway cohesion (damped)", 1),
         ("ppi 0.20 — STRING network proximity (damped)", 1),
         ("clinical 0.15 — trials + literature for THIS pair", 1),
         ("indication 0.10 — existing-use token match", 1),
         ("regulatory 0.10 — development phase / orphan", 1)],
        [("Bounded bonuses (additive)", 0, True),
         ("+ HetioNet network evidence (Compound→Gene←Disease)", 1),
         ("+ directional literature (typed drug→gene→disease, DRKG)", 1),
         ("Consistency", 0, True),
         ("One function, every surface → same number everywhere", 1),
         ("Full breakdown returned for transparency", 1)],
    ),
    footer="Shared by discover, repurpose, pathway, novel-targets and the dossier")

content_slide(
    "Applicable-Weight Renormalization", "The 2026-07 fix for systematic harshness",
    [
        ("Problem: for a genuinely NOVEL pair, indication≈0 and clinical sits at its floor → ~35% of weight is dead", 0, True),
        ("Result was a hard ceiling (~0.10) — real leads read as a '10 / 100' regardless of mechanism strength", 0, True),
        ("Fix: renormalize the disease-specific mechanistic mass over only the dimensions that carry signal", 0, True),
        ("Regulatory kept as a small GLOBAL prior — NOT amplified (avoids rewarding an unrelated approved drug)", 0, True),
        ("Denominator floor prevents a lone weak signal from being over-amplified — separates leads, not inflates", 0, False),
        ("Validated: Metformin→Alzheimer 0.36→0.47; Riluzole→ALS 0.47→0.55; weak/random pairs stay at the floor", 0, False),
    ],
    "A candid engineering slide. The original weighted mean had a structural flaw: for a novel "
    "repurposing pair, two of the six dimensions are near-zero BY CONSTRUCTION — if the drug were "
    "already indicated, it wouldn't be repurposing — so about a third of the weight was dead weight, "
    "capping even strong mechanistic leads at a low number. The fix renormalizes the mechanistic "
    "mass over the dimensions that actually carry disease-specific signal, while keeping the global "
    "regulatory phase as a small, non-amplified prior so an approved-but-mechanistically-irrelevant "
    "drug can't inflate. A denominator floor stops a single weak signal from running away. The result "
    "improves DISCRIMINATION — genuine leads rise, random pairs stay low — which we verified on known "
    "examples.")

content_slide(
    "Evidence Sub-Signals", "How each dimension is computed",
    None,
    "The supporting modules behind the score. Hub-degree damping (Himmelstein's DWPC idea) down-weights "
    "promiscuous hub genes like TP53/EGFR so they can't max out cohesion for every disease. The "
    "signature engine adds DIRECTION — does the drug REVERSE the disease's expression signature "
    "(favourable) or MIMIC it (unfavourable) — a CMap/LINCS-style connectivity idea applied to "
    "pathway scoring. Directional evidence uses typed drug→gene→disease triples from DRKG (real "
    "relational support, not mere co-occurrence). Together these make the mechanistic score "
    "direction-aware and hub-robust rather than a naive overlap count.",
    two_col=(
        [("Mechanistic", 0, True),
         ("hub_degree.py — DWPC damping of network hubs", 1),
         ("signature_engine.py — signed connectivity (reverse vs mimic disease signature)", 1),
         ("target_coverage.py — polygenic driver-set coverage", 1)],
        [("Literature & network", 0, True),
         ("directional_evidence.py — typed DRKG triples (drug→gene→disease)", 1),
         ("repurposing_scorer.py — HetioNet Compound→Gene←Disease paths", 1),
         ("kg_embedding.py / kg_score.py — DistMult link-prediction signal", 1)],
    ))

content_slide(
    "Guardrails — The Penalty Layer", "Down-ranking implausible or dead pairs (honestly)",
    None,
    "Guardrails are multiplicative penalties that encode clinical and regulatory reality. The safety "
    "cross-filter uses FAERS disproportionality to penalise a drug whose serious toxicity hits the "
    "disease's organ system. The Clinical Constraint Harmonizer penalises a toxic drug used for a "
    "benign or chronic indication (severity/duration mismatch). Target-coverage penalises a drug that "
    "hits only part of a polygenic driver set. Two CTPA rules kill phantoms: a pair with no functional "
    "cohesion beyond a single string match, and a program that already failed/was abandoned in that "
    "indication (registry ghost). Crucially, in 2026 we COLLAPSED the soft penalties (take the single "
    "worst instead of multiplying) so they stop double-counting and crushing every score — while hard, "
    "evidence-based kills still stack.",
    two_col=(
        [("Soft penalties (take the worst, not the product)", 0, True),
         ("safety_filter.py — FAERS toxicity ↔ disease-organ", 1),
         ("clinical_constraints.py — severity / chronicity fit (CCH)", 1),
         ("target_coverage.py — incomplete polygenic coverage", 1)],
        [("Hard gates (evidence-based, still stack)", 0, True),
         ("CTPA Rule 1 — phantom off-target (no cohesion)", 1),
         ("CTPA Rule 2 — registry ghost (registry_audit.py)", 1),
         ("Trial-failed-here — a real terminated trial in this pair", 1),
         ("2026 fix: collapse soft penalties → stop double-counting", 1)],
    ))

content_slide(
    "Reverse Repurposing Engine", "reverse_repurposing.py — drug → new indications (505(b)(2))",
    [
        ("Enter a molecule you already hold approval for → derive targets → find NEW candidate indications", 0, True),
        ("Candidates from three orthogonal sources: Open Targets genetics + ClinicalTrials.gov + PubMed co-mention", 0, True),
        ("Excludes on-label indications up front (a new use, not a confirmation)", 0, True),
        ("Ranks by the canonical pair score; reports 505(b)(2) feasibility separately (not in the rank)", 0, True),
        ("Phenotype filter: drops toxicity phenotypes/signs that aren't repurposable diseases", 0, False),
        ("This is the Alembic engagement's core surface (/repurpose)", 0, False),
    ],
    "The commercial centrepiece. Reverse repurposing takes a drug the client already owns and finds "
    "new indications for it. Candidates come from three orthogonal signals so we don't miss real-world "
    "leads that pure genetics would: genetic associations, existing trials of the drug (which reveal "
    "human interest that genetics can't), and literature co-mention. We strip out on-label uses so the "
    "output is genuinely new. Ranking uses the shared canonical score; the 505(b)(2) feasibility and "
    "regulatory status are reported per candidate but deliberately kept OUT of the ranking because a "
    "drug's global status is constant and doesn't discriminate between indications. This is the "
    "/repurpose page the Alembic work centres on.")

content_slide(
    "Pathway-First Screen", "pathway_screen.py — enter through biology",
    [
        ("Enter a biological PATHWAY (e.g. 'mTOR signaling') and a desired direction (suppress / activate)", 0, True),
        ("Pathway → member genes → drugs that modulate them (direction-aware) + diseases the pathway drives (Open Targets)", 0, True),
        ("Emits genuine drug↔indication repurposing pairs, scored with the SAME canonical pair score", 0, True),
        ("Mode B (discovery): pathway → inferred diseases; Mode A (anchored): pathway as a filter on a disease screen", 0, True),
        ("Fuzzy pathway name resolution (Greek letters, spelling variants: NF-kB ↔ NF-kappaB)", 0, False),
        ("Modulation score = direction-match + coverage + phase; fed to the scorer as a mechanistic prior", 0, False),
    ],
    "The pathway screen lets a biologist start from a mechanism rather than a disease. You give it a "
    "pathway and the direction you want (suppress or activate); it finds the drugs that modulate the "
    "pathway's genes in that direction and the diseases the pathway drives, then emits real drug-to-new"
    "-indication pairs scored by the shared engine. A subtle but important point: the pathway "
    "modulation the screen establishes is passed INTO the canonical scorer as a mechanistic prior, so "
    "a candidate isn't scored as if it had zero mechanism when the very reason it surfaced IS a "
    "mechanism. Name resolution is fuzzy so spelling and Greek-letter variants all resolve.")

content_slide(
    "Novel-Target Discovery", "novel_targets.py — from validation to discovery",
    [
        ("Infers targets NOT yet annotated to a disease via PPI guilt-by-association off the known seed genes", 0, True),
        ("disease → OT seed genes (genetic weight) → STRING PPI partners → a well-connected non-seed gene = candidate", 0, True),
        ("Each inferred target carries its evidence (which seeds, at what confidence) + a normalized confidence", 0, True),
        ("drugs_via_novel_targets() closes the loop: drugs hitting inferred targets = leads reachable ONLY via novelty", 0, True),
        ("Clearly labelled INFERRED (network-based), not an experimentally annotated association — fail-soft", 0, False),
    ],
    "Where the other engines confirm or extend known biology, novel-target discovery tries to find "
    "NEW biology. It takes the disease's known genetic seed genes, expands along the STRING "
    "protein-protein interaction network, and flags non-seed genes that are heavily wired into the "
    "seed set — hidden nodes one hop off the known disease module. Then it finds drugs that hit those "
    "inferred targets, giving repurposing leads reachable only through the novel target. Everything is "
    "explicitly labelled as network-inferred, not experimentally annotated, and it fails soft to an "
    "empty result rather than fabricating a target.")

content_slide(
    "Positive Pivots & Class Comparison", "Turning a mismatch into a lead",
    [
        ("positive_pivots.py — when a pair is penalised (e.g. too toxic), suggest HOW to make it work", 0, True),
        ("  • pivot to a high-severity / orphan variant where the therapeutic window opens", 1),
        ("  • pivot to a dose-sparing combination with a low-toxicity partner on the same target", 1),
        ("class_comparison.py — within-class vs out-of-class differentiation (is the effect class-wide or drug-specific?)", 0, True),
        ("approval_supplement.py — catches ChEMBL's per-indication phase lag so status is current", 0, False),
    ],
    "Two modules that add strategic nuance. Positive Pivots is a nice touch: rather than just "
    "penalising a hard pair and moving on, it proposes concrete hypotheses to rescue it — retarget a "
    "severe or orphan variant of the disease where a toxic drug is justified, or combine an ultra-low "
    "dose with a low-toxicity partner acting on the same target. Class comparison asks whether an "
    "observed effect is class-wide (a target effect) or specific to this molecule, which matters for "
    "IP and differentiation. The approval supplement corrects ChEMBL's tendency to lag on "
    "per-indication phase.")

section_slide("03", "Calibration, Validation & the Predictor",
    "Section divider. This is the scientific-integrity heart of the platform: how we make a raw "
    "score interpretable, how we validated it against a gold standard, what that validation "
    "revealed, and the real predictor we built as a result.")

content_slide(
    "Score Calibration", "score_calibration.py + build_score_null.py",
    [
        ("A raw composite is meaningless without a baseline — a novel lead legitimately scores low in absolute terms", 0, True),
        ("build_score_null.py — scores RANDOM drug-disease pairs (from RepoDB names) → a null/background distribution", 0, True),
        ("Reports ENRICHMENT — how many× above the random-pair median a pair scores (discriminating, honest)", 0, True),
        ("Tier bands (Strong / Promising / Moderate / Weak) — currently from raw-score bands", 0, True),
        ("Percentile saturates once signal separates from noise → we lead with enrichment, not percentile", 0, False),
    ],
    "Calibration answers 'is this number good?' A raw score has no meaning without a reference, so we "
    "built a null distribution by scoring random drug-disease pairs. A real pair's ENRICHMENT — how "
    "far above the random-pair median it sits — is a defensible, discriminating statement (e.g. '9× "
    "above background'). We report tiers from raw-score bands. This slide sets up the honest question "
    "the client rightly pushed on next: those tier bands were hand-chosen. That question drove the "
    "RepoDB validation — the next three slides.")

content_slide(
    "The RepoDB Validation", "Does the score separate SUCCESS from FAILURE?",
    [
        ("RepoDB = 13,558 labelled pairs (8,931 Approved / 4,627 Failed) — the outcome gold standard", 0, True),
        ("repodb_validate.py — scored a stratified sample of Approved vs Failed with the hand-built composite", 0, True),
        ("Result: composite AUC 0.41 (INVERTED — ranks failures above successes); mechanistic-only 0.54", 0, True),
        ("Why: trial activity ≠ success — a Phase-3 FAILURE is a heavily-studied, well-connected pair", 0, True),
        ("Honest finding: the hand-built score is NOT a validated success predictor — cutoffs from it are arbitrary", 0, False),
    ],
    "The client asked exactly the right question: how do we decide the cutoff number, and did we use "
    "RepoDB? We built a validation harness and got an uncomfortable, important answer. Scoring real "
    "Approved-versus-Failed pairs, the hand-built composite achieved an AUC of 0.41 — BELOW random — "
    "meaning it actually ranked failures above successes. The reason is deep: a drug that reached a "
    "big Phase-3 trial and failed is a heavily studied, densely connected pair, so any "
    "connectivity-based score rewards it. The honest conclusion: the composite measures plausibility, "
    "not success, and any 'probability of success' cutoff derived from it would be meaningless. That "
    "finding forced the next step — a real, trained predictor.")

content_slide(
    "The DWPC Plausibility Predictor", "metapath_features.py + train_repurposing_predictor.py",
    [
        ("Leakage-free DWPC / metapath features (Rephetio; Himmelstein 2017) — typed multi-hop paths, hub-damped", 0, True),
        ("No metapath traverses the 'treats' edge → predicting treatment is leakage-free by construction", 0, True),
        ("Trained a logistic model on the Hetionet gold standard (755 Compound-treats-Disease edges)", 0, True),
        ("Held-out AUC 0.978 — and 0.978 COMPOUND-DISJOINT (generalises to drugs it never trained on)", 0, True),
        ("Data-derived actionable cutoff: P(treats) ≥ 0.15 (Youden-optimal) — a real number, not a hand-set band", 0, True),
    ],
    "Instead of hand-set bands, we trained a real model using the DWPC / metapath approach from the "
    "Rephetio paper — the method that actually works on Hetionet. It counts typed multi-hop paths "
    "between a drug and a disease, down-weighting hub nodes, and crucially none of the paths use the "
    "'treats' edge, so predicting treatment from them is leakage-free. Trained on the 755 "
    "gold-standard treats edges, it reaches an AUC of 0.978 — and, importantly, the same 0.978 when we "
    "hold out ENTIRE drugs, proving it generalises to molecules it never saw rather than memorising. "
    "From its out-of-fold predictions we derive a real, Youden-optimal cutoff of 0.15. This is the "
    "rigorous answer to 'how do we decide a number.'")

content_slide(
    "Plausibility ≠ Success — the honest scope", "repodb_external_validate.py + repurposing_predictor.py",
    [
        ("External test: applied the AUC-0.98 predictor to RepoDB Approved-vs-Failed (hard, real failures)", 0, True),
        ("Result: AUC 0.42 — it ranks plausibility above RANDOM, but NOT successes above failures", 0, True),
        ("Fundamental ceiling: efficacy / safety / PK are not in any knowledge graph — both classes are 'plausible'", 0, True),
        ("So we ship it HONESTLY: a mechanistic-PLAUSIBILITY / triage score, explicitly NOT a success probability", 0, True),
        ("repurposing_predictor.py serves P(treats) where the pair maps into Hetionet; fail-soft otherwise", 0, True),
        ("UI badge shows BOTH AUCs (0.98 vs random, 0.42 vs failures) — it can never be mistaken for an oracle", 0, False),
    ],
    "The most important integrity slide. We stress-tested the excellent predictor against RepoDB's "
    "REAL failures — the hard test — and it scored 0.42: it separates plausible from random, but "
    "cannot separate eventual successes from eventual failures, because both are biologically "
    "plausible and the deciding factors — efficacy, safety, dose, pharmacokinetics — simply aren't in "
    "the graph. This is a ceiling no cheap model escapes. So we ship the predictor for what it "
    "genuinely is: a validated mechanistic-plausibility triage number with a real cutoff, wired in "
    "additively and fail-soft, and labelled in the UI with BOTH AUCs so it can never be oversold as a "
    "probability of clinical success. That intellectual honesty is a feature, not an apology.")

content_slide(
    "Probability-of-Success (PoS) Model", "pos_model.py + train_pos_model.py",
    [
        ("Analytic phase-transition model: likelihood of progressing through each remaining clinical phase", 0, True),
        ("Repurposing-aware: an approved molecule entering a new indication has Phase-1 safety largely de-risked", 0, True),
        ("Optional fitted calibrator (biological-profile features) — blended 50/50 ONLY when trustworthy", 0, True),
        ("Leakage guard: refuses to certify a calibrator whose AUC is implausibly high (reverse-causation)", 0, True),
        ("Shown per candidate as the asset-scoring 'probability of success' on the 505(b)(2) dossier", 0, False),
    ],
    "Separate from the plausibility predictor, the PoS model estimates the probability of progressing "
    "through the remaining clinical phases, using published phase-transition base rates and adjusting "
    "for the fact that a repurposed approved molecule has much of Phase-1 safety already cleared. "
    "There is an optional machine-learned calibrator, but it is only blended in when it passes a "
    "leakage integrity gate — we explicitly REFUSE a model whose accuracy is suspiciously high, "
    "because that signals it's memorising the label via reverse causation (annotation richness "
    "reflecting approval). Another example of the honesty posture: we would rather use a modest "
    "analytic model than a leaky, impressive-looking one.")

content_slide(
    "Knowledge Graph", "biocypher_graph.py + kg_embedding.py + kg_score.py",
    [
        ("biocypher_graph.py — assembles the story graph: the evidence chain behind a hypothesis (drug→gene→pathway→disease)", 0, True),
        ("Hetionet v1.0 in the database — the curated, typed biomedical graph (nodes + metaedges)", 0, True),
        ("kg_embedding.py — DistMult knowledge-graph embeddings for link prediction (data/kg_embeddings.npz)", 0, True),
        ("kg_score.py — embedding-based repurposing score + Hetionet crosswalk", 0, True),
        ("Knowledge Graph page (/graph) — interactive evidence-chain visualisation, molecule- and disease-aware", 0, False),
    ],
    "The knowledge-graph layer serves two purposes. Explanatory: the BioCypher story graph turns a "
    "score into a visible evidence chain — this drug hits this gene, which sits in this pathway, which "
    "the disease implicates — which is what makes a hypothesis credible to a scientist. Predictive: we "
    "hold DistMult embeddings of the graph for link prediction, and the DWPC metapath features (from "
    "the predictor section) are computed on the same Hetionet graph. The /graph page renders the "
    "evidence chain interactively and adapts to whether the user searched a molecule or a disease.")

section_slide("04", "Structure, Chemistry & Pharmacokinetics",
    "Section divider. The physical-science stack: real structure-based docking, quantum chemistry, "
    "PBPK pharmacokinetics, and the drug-like property scorers. This is what lets the platform go "
    "beyond associations to physical plausibility.")

content_slide(
    "Molecular Docking — Architecture", "docking_service.py — a priority chain",
    [
        ("Priority chain: real AutoDock Vina (default) → optional DiffDock / Boltz-2 → physics fallback → descriptor estimate", 0, True),
        ("real_pdb_fetcher.py — gene → UniProt identity → best RCSB structure (prefers HOLO) → else AlphaFold (pLDDT-filtered)", 0, True),
        ("pocket_finder.py — detects, ranks and NAMES binding pockets; prefers the co-crystal (functional) site", 0, True),
        ("Pose validation — clash + pocket-containment checks; drops floating/clashing poses", 0, True),
        ("Honest 'no structure' path — never dresses a placeholder peptide up as a real docking result", 0, False),
    ],
    "Docking is a priority chain that always returns something honest. It fetches a REAL protein "
    "structure by resolving the gene to a UniProt accession and picking the best RCSB experimental "
    "structure — now preferring HOLO structures that have a bound ligand, so there's a real binding "
    "site — falling back to a pLDDT-filtered AlphaFold model. The pocket finder detects and NAMES "
    "pockets and prefers the co-crystal functional site over the largest random cavity, which was the "
    "key fix for 'docking into the wrong place.' Poses are validated for clashes and containment. If "
    "no real structure exists, the platform says so rather than faking one.")

content_slide(
    "Real AutoDock Vina — Native Windows", "vina_engine.py + vendor/vina_worker.py",
    [
        ("Replaced the in-house 'Vina-style' approximation with REAL AutoDock Vina 1.2.7", 0, True),
        ("Windows reality: conda vina/fpocket have no win-64 build → use the official Vina Windows CLI binary", 0, True),
        ("Worker chain (drugdisc-agent env): Meeko ligand PDBQT → pdbfixer + OpenBabel receptor → box on the co-crystal site → Vina", 0, True),
        ("App shells out to the worker; poses returned as SDF + real ΔG (kcal/mol); fail-soft to the physics engine", 0, True),
        ("VALIDATED: imatinib redocked into 1IEP = −12.6 kcal/mol, 0.24 Å from the crystal ligand; EGFR/erlotinib −8.65", 0, True),
    ],
    "The docking upgrade this cycle. We replaced a triage-grade in-house approximation with the real "
    "AutoDock Vina engine. Because Vina's Python bindings have no Windows conda build, we use the "
    "official Vina Windows binary and run the preparation chemistry — Meeko for the ligand, pdbfixer "
    "and OpenBabel for the receptor — in the isolated drugdisc-agent environment, with the docking box "
    "centred on the real co-crystal binding site. The app just shells out and reads back SDF poses and "
    "a true binding energy. We validated it properly: redocking imatinib into its own crystal "
    "reproduced the pose to within a quarter of an Ångström at a realistic −12.6 kcal/mol. This is "
    "real structure-based science, not a heuristic.")

content_slide(
    "Structure Modelling Options", "pocket_finder.py + boltz_engine.py + dock_engine.py",
    [
        ("boltz_engine.py — Boltz-2: structure + BINDING-AFFINITY prediction (opt-in; GPU/weights) for a predicted ΔG", 0, True),
        ("dock_engine.py — the always-available physics fallback: RDKit conformers + Monte-Carlo pose + Vina-style scoring", 0, True),
        ("Consensus: externally-generated poses (DiffDock/Vina) re-scored by the empirical function as an independent check", 0, True),
        ("AtomisticSkills recipes (vendored) — protein-prep, ligand-prep, binding-site, redocking-RMSD skills", 0, False),
        ("DiffDock (opt-in) — deep-learning blind docking when the user wants it", 0, False),
    ],
    "Beyond Vina, the platform offers alternatives. Boltz-2 predicts structure AND binding affinity "
    "directly (opt-in, needs GPU). The built-in physics engine is the guaranteed fallback — it "
    "actually translates and rotates RDKit conformers into the pocket and scores them with a Vina-"
    "style empirical function — so docking never fails to return a spatial pose. When an external "
    "engine produces poses, we cross-score them with the empirical function as an independent "
    "consensus check. The AtomisticSkills library provides the vetted recipes for protein and ligand "
    "prep and for redocking-RMSD validation.")

content_slide(
    "Quantum Chemistry Engine", "qc_engine.py — GFN2-xTB",
    [
        ("REAL semi-empirical quantum chemistry (GFN2-xTB) via a dedicated micromamba env at .qc/", 0, True),
        ("Geometry optimization + electronic properties: HOMO/LUMO gap, dipole, solvation ΔG (implicit water)", 0, True),
        ("Surfaces on the Analysis page 'Optimization' tab — molecular-level physical descriptors", 0, True),
        ("Isolated env → the app never imports the QM stack; results returned as structured JSON", 0, False),
    ],
    "The quantum-chemistry engine runs genuine GFN2-xTB semi-empirical calculations — geometry "
    "optimization and electronic-structure properties like the HOMO-LUMO gap, dipole moment and "
    "implicit-solvent free energy. It runs in its own micromamba environment so its heavy dependencies "
    "stay isolated, and the results appear on the Analysis page's optimization tab. This is another "
    "'real, not simulated' component — the numbers come from an actual quantum-chemistry code, not a "
    "lookup or an approximation.")

content_slide(
    "PBPK Pharmacokinetics", "pbpk_simulation.py",
    [
        ("Real perfusion-limited PBPK model — a system of ODEs solved with scipy LSODA", 0, True),
        ("Predicts tissue/plasma concentration-time profiles from dose, physiology and drug properties", 0, True),
        ("Free-drug target occupancy vs a MEASURED ChEMBL Ki — connects PK to pharmacology", 0, True),
        ("Answers: does a repurposing dose actually reach and engage the target at the disease site?", 0, False),
    ],
    "The PBPK module is a real physiologically-based pharmacokinetic model — a perfusion-limited "
    "compartment system solved as ODEs with scipy's LSODA integrator. It predicts how drug "
    "concentration evolves in tissues and plasma, and then compares the free-drug concentration "
    "against a MEASURED binding affinity from ChEMBL to estimate target occupancy. That closes an "
    "important loop for repurposing: it's not enough for a drug to hit the right target in principle — "
    "it has to actually reach the disease tissue at a concentration that engages the target at a "
    "feasible dose.")

content_slide(
    "Developability & CNS Exposure", "developability.py + cns_mpo.py",
    [
        ("developability.py — area-aware drug-likeness: oral (Lipinski + Veber), CNS (BBB), ophthalmic (corneal)", 0, True),
        ("Scores a molecule against the RIGHT rule set for the searched disease's route of administration", 0, True),
        ("cns_mpo.py — CNS-MPO (Wager 2010) multiparameter score for brain penetration (0–6)", 0, True),
        ("Feeds a BBB-risk flag on CNS indications and the CCH tissue-barrier factor", 0, False),
    ],
    "Two drug-property scorers. Developability is area-aware — it applies the correct rule set for the "
    "route implied by the disease: Lipinski and Veber for oral systemic drugs, blood-brain-barrier "
    "criteria for CNS indications, corneal-permeation criteria for eye drugs. CNS-MPO is the published "
    "Wager multiparameter score for central-nervous-system exposure, which drives a blood-brain-"
    "barrier risk flag and feeds the clinical-constraint tissue factor. Together they answer 'is this "
    "molecule actually developable for THIS indication's delivery route.'")

section_slide("05", "Decision Context, Regulatory & Delivery",
    "Section divider. The last layer wraps the science in the commercial and regulatory context a "
    "505(b)(2) decision needs: the evidence dossier, regulatory/IP status, portfolio fit, and the "
    "cross-cutting infrastructure that holds it all together.")

content_slide(
    "Evidence Dossier & Regulatory Verdict", "evidence_dossier.py + regulatory_verdict.py",
    [
        ("evidence_dossier.py — assembles the per-candidate 505(b)(2) dossier: evidence chain, scores, PoS, structure, PK", 0, True),
        ("regulatory_verdict.py — single source of truth for novelty + regulatory status (phase-aware)", 0, True),
        ("Novelty logic: only an APPROVED-here molecule disqualifies a 505(b)(2) claim; prior lower-phase work de-risks it", 0, True),
        ("orange_book.py — authoritative FDA patent & exclusivity status (is the IP window open?)", 0, True),
        ("approval_supplement.py — corrects per-indication phase lag so the verdict is current", 0, False),
    ],
    "The dossier is where a hypothesis becomes a decision document. It assembles the full evidence "
    "chain, the scores and their breakdown, the probability of success, the docking and PK results "
    "into a single 505(b)(2) write-up. The regulatory verdict is the single source of truth for "
    "novelty and status, and it's phase-aware: the only thing that truly disqualifies a 505(b)(2) "
    "claim is the molecule already being APPROVED for that exact disease — prior lower-phase "
    "development actually DE-RISKS the asset rather than blocking it. The Orange Book integration tells "
    "you whether the patent and exclusivity window is open, which is decisive for commercial viability.")

content_slide(
    "Global Regulatory Landscape & Signals", "global_landscape.py + acquisition_signal.py",
    [
        ("global_landscape.py — one consolidated US / EU / India regulatory view (real data, not link tiles)", 0, True),
        ("Pulls FDA, EMA, MHRA, CTRI/CDSCO to show where a drug↔indication already stands globally", 0, True),
        ("acquisition_signal.py — asset-acquisition signal from Orange Book / market status (is this a buy signal?)", 0, True),
        ("class_comparison.py — is the effect class-wide or drug-specific (differentiation / freedom-to-operate)?", 0, False),
    ],
    "Repurposing decisions are global, so the landscape module consolidates US, EU and Indian "
    "regulatory status into one data-driven view rather than a page of external links — where does "
    "this drug-indication pair already stand with FDA, EMA, MHRA and the Indian authorities. The "
    "acquisition-signal engine reads Orange Book and market status to flag assets that look like "
    "acquisition or in-licensing opportunities. Class comparison informs differentiation and "
    "freedom-to-operate. This is the commercial intelligence layer on top of the science.")

content_slide(
    "Portfolio Matching", "portfolio.py — fit to the client's assets",
    [
        ("Matches surfaced candidates against the client's existing portfolio (the Alembic 505(b)(2) engagement)", 0, True),
        ("Flags and boosts molecules already in-house → immediately actionable, 505(b)(2)-ready leads", 0, True),
        ("Turns a generic repurposing hit into a specific, ownable opportunity for the partner", 0, False),
    ],
    "Portfolio matching makes the output specific to the client. It cross-references the surfaced "
    "candidates against the partner's existing asset portfolio and flags — and boosts — molecules "
    "they already own, because those are the immediately actionable 505(b)(2) opportunities: no "
    "in-licensing needed, existing manufacturing, a clear path. For the Alembic engagement this is "
    "what turns a generic repurposing hit into 'here is a molecule you already hold that could open a "
    "new indication.'")

content_slide(
    "Cross-Cutting Infrastructure", "The glue that makes it robust",
    [
        ("config.py — single source of DB and environment configuration (no hardcoded credentials; .env-driven)", 0, True),
        ("http_client.py — one HTTP client for every external call: retry, backoff, sane timeouts", 0, True),
        ("diskcache.py — disk-backed cache shared across services (30-day TTLs for structures, etc.)", 0, True),
        ("safety_filter.py — shared FAERS toxicity profiles reused by the scorer and CCH (one network, cached)", 0, True),
        ("Fail-soft everywhere — a down dependency degrades a feature, never the platform", 0, True),
    ],
    "The unglamorous but essential layer. One config module is the single source of truth for "
    "database and environment settings, with no hardcoded credentials — everything comes from the "
    "environment. One HTTP client wraps every external call with retry, backoff and sensible timeouts. "
    "One disk cache is shared across services with long TTLs so we don't re-fetch structures or "
    "associations. The FAERS safety profiles are computed once and reused by both the scorer and the "
    "clinical-constraint layer. And the pervasive design rule is fail-soft: if any dependency is "
    "unavailable, the affected feature degrades gracefully while the rest of the platform keeps "
    "working. This robustness is what makes it demo-able and trustworthy.")

section_slide("06", "Validation Posture & Roadmap",
    "Final section. Pull the honesty theme together into a validation scorecard, and lay out what's "
    "next.")

content_slide(
    "Data-Quality & Validation Posture", "Built for a skeptical pharma audience",
    [
        ("Pharma distrust is of the DATA — so we validate once, at the data layer (the DATA_QUALITY workstream)", 0, True),
        ("ChEMBL 33 restore-verified; POZ + cipherq share one integrity-checked database; a restore-verified backup exists", 0, True),
        ("Every score has a transparent breakdown; every penalty is flagged with its reason", 0, True),
        ("Confidence is stated honestly — enrichment vs random, validated plausibility where covered, PoS separately", 0, True),
        ("The Validation page (/validation) tells the data-quality story end-to-end", 0, False),
    ],
    "The platform's design centre of gravity is trust. Pharma skepticism is fundamentally about the "
    "DATA, so we validate it once at the data layer rather than per query: ChEMBL 33 is "
    "restore-verified and integrity-checked, shared cleanly across projects, with a verified backup. "
    "Every score exposes its full breakdown and every penalty states its reason, so nothing is a black "
    "box. And confidence is always stated honestly — enrichment against a random background, the "
    "validated plausibility probability where the pair is in the graph, and probability of success "
    "reported as its own separate number. The Validation page narrates this end to end.")

content_slide(
    "Honesty Scorecard", "What is validated vs. what is an estimate",
    None,
    "The single most important slide for a pharma partner. Be explicit about epistemic status. "
    "VALIDATED against a gold standard: the docking (redocking RMSD), the DWPC plausibility predictor "
    "(held-out AUC 0.98, compound-disjoint). REAL computation but not outcome-validated: the QM "
    "engine, PBPK, the mechanistic evidence scores — these are genuine physics/biology, and useful, "
    "but they don't claim to predict clinical success. HEURISTIC / estimate: the composite tier bands, "
    "the PoS analytic base rates, developability rules — reasonable, transparent, but not fitted. And "
    "the explicit NEGATIVE result: clinical success is NOT predictable from mechanism/network (AUC "
    "0.42), and we say so rather than pretend. Leading with this scorecard is what earns trust.",
    two_col=(
        [("Validated (gold standard)", 0, True),
         ("DWPC plausibility predictor — AUC 0.978, compound-disjoint", 1),
         ("Docking — imatinib redock 0.24 Å, −12.6 kcal/mol", 1),
         ("Real computation (not outcome-validated)", 0, True),
         ("Quantum chemistry (GFN2-xTB), PBPK ODE model", 1),
         ("Mechanistic evidence scores (target/pathway/PPI)", 1)],
        [("Heuristic / transparent estimate", 0, True),
         ("Composite tier bands; PoS analytic base rates", 1),
         ("Developability / CNS-MPO rule sets", 1),
         ("Stated NEGATIVE result (honesty)", 0, True),
         ("Clinical SUCCESS not predictable from graph (AUC 0.42)", 1),
         ("→ plausibility labelled as triage, never as an oracle", 1)],
    ))

content_slide(
    "Roadmap", "Where it goes next",
    [
        ("Refresh the knowledge graph beyond Hetionet (~2016) → wider, current coverage for the predictor", 0, True),
        ("Explore a weak-but-real SUCCESS signal: genetic-evidence + clinical features (Nelson 2015 ~2× effect)", 0, True),
        ("Clean pose chemistry (Meeko RDKitMolCreate export) + automated redocking-RMSD gate on every dock", 0, True),
        ("Broaden portfolio & regulatory coverage (more registries, live IP monitoring)", 0, True),
        ("Productionize compute (job queue for docking/QM; precomputed DWPC cache shipped)", 0, False),
    ],
    "Where we take it next. The biggest lever is a fresher, broader knowledge graph — Hetionet is a "
    "2016 snapshot, which limits coverage and misses newer drugs, so a modern graph would widen the "
    "predictor's reach. A genuine (if modest) success signal may be reachable by adding "
    "genetic-evidence and clinical features — genetically supported targets are about twice as likely "
    "to be approved — but we'd validate it honestly before claiming it. On the engineering side: "
    "cleaner pose chemistry with an automated redocking-RMSD gate, broader regulatory and IP coverage, "
    "and productionizing the heavy compute behind a job queue.")

title_slide(
    "SUMMARY",
    "One engine, honestly measured",
    "Disease-agnostic repurposing across forward, reverse, pathway and novel-target entries — scored "
    "on real mechanistic, clinical, structural and regulatory evidence, with a validated plausibility "
    "predictor and an explicit, honest boundary on what it can and cannot claim.",
    "Close by restating the two things that make RepurposeIQ credible. First, breadth done right: one "
    "shared, consistent scoring core serves every entry point — forward discovery, the reverse "
    "505(b)(2) search, pathway-first and novel-target — layered with real structure-based docking, "
    "quantum chemistry, PBPK and a knowledge-graph predictor, all disease-agnostic and fail-soft. "
    "Second, and more important for a pharma partner: intellectual honesty. We validated against a "
    "gold standard, we found and shipped a real predictor, and we drew a clear line — plausibility we "
    "can measure and stand behind; clinical success we cannot predict from mechanism and we say so. "
    "That combination — real science plus honest confidence — is the platform's actual moat. Invite "
    "questions.")

prs.save("RepurposeIQ_Platform_Architecture.pptx")
print(f"Saved RepurposeIQ_Platform_Architecture.pptx — {len(prs.slides.__iter__.__self__._sldIdLst)} slides")
